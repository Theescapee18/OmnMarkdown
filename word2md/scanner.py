"""
OmnMarkdown 敏感词扫描模块
- 自定义敏感词库（本地 JSON 存储）
- 批量扫描 .docx / .md / .txt 等文件
- 生成扫描报告（JSON + 终端输出）
"""

import json
import re
import os
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional, Tuple


# 默认词库路径（与可执行文件同级目录）
DEFAULT_WORDLIST = Path.home() / ".omnmarkdown" / "sensitive_words.json"
DEFAULT_REPORT_DIR = Path.home() / ".omnmarkdown" / "scan_reports"


class SensitiveWordScanner:
    """敏感词扫描器"""

    def __init__(
        self,
        wordlist_path: Optional[str] = None,
        report_dir: Optional[str] = None,
    ):
        self.wordlist_path = Path(wordlist_path) if wordlist_path else DEFAULT_WORDLIST
        self.report_dir = Path(report_dir) if report_dir else DEFAULT_REPORT_DIR
        self.words: List[str] = []
        self.categories: Dict[str, List[str]] = {}
        self._load_wordlist()

    # ================================================================
    # 词库管理
    # ================================================================

    def _load_wordlist(self):
        """加载敏感词库"""
        if not self.wordlist_path.exists():
            self.wordlist_path.parent.mkdir(parents=True, exist_ok=True)
            self._save_wordlist()
            return

        try:
            data = json.loads(self.wordlist_path.read_text(encoding="utf-8"))
            if isinstance(data, list):
                self.words = data
                self.categories = {"default": data}
            elif isinstance(data, dict):
                self.categories = data
                self.words = []
                for cat_words in data.values():
                    self.words.extend(cat_words)
            self.words = list(dict.fromkeys(self.words))  # 去重保序
        except Exception:
            self.words = []
            self.categories = {}

    def _save_wordlist(self):
        """保存敏感词库"""
        self.wordlist_path.parent.mkdir(parents=True, exist_ok=True)
        data = self.categories if self.categories else {"default": self.words}
        self.wordlist_path.write_text(
            json.dumps(data, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )

    def add_words(self, words: List[str], category: str = "default") -> int:
        """添加敏感词，返回新增数量"""
        added = 0
        if category not in self.categories:
            self.categories[category] = []

        for w in words:
            w = w.strip()
            if w and w not in self.words:
                self.words.append(w)
                self.categories[category].append(w)
                added += 1

        if added > 0:
            self._save_wordlist()
        return added

    def remove_words(self, words: List[str]) -> int:
        """删除敏感词，返回删除数量"""
        removed = 0
        for w in words:
            w = w.strip()
            if w in self.words:
                self.words.remove(w)
                removed += 1
                for cat_words in self.categories.values():
                    if w in cat_words:
                        cat_words.remove(w)

        if removed > 0:
            self._save_wordlist()
        return removed

    def list_words(self) -> Dict[str, List[str]]:
        """列出所有分类及其敏感词"""
        return dict(self.categories)

    def get_word_count(self) -> int:
        return len(self.words)

    # ================================================================
    # 文件扫描
    # ================================================================

    def scan_file(self, file_path: str) -> List[Dict]:
        """扫描单个文件，返回命中记录列表"""
        file_path = Path(file_path)
        ext = file_path.suffix.lower()
        hits = []

        if not file_path.exists():
            return hits

        if not self.words:
            return hits

        # 提取文本内容
        text = self._extract_text(file_path, ext)
        if not text:
            return hits

        # 逐行扫描
        lines = text.split("\n")
        for line_num, line in enumerate(lines, 1):
            for word in self.words:
                # 不区分大小写匹配
                pattern = re.escape(word)
                for match in re.finditer(pattern, line, re.IGNORECASE):
                    # 提取上下文（前后各 20 字符）
                    start = max(0, match.start() - 20)
                    end = min(len(line), match.end() + 20)
                    context = line[start:end].strip()
                    if start > 0:
                        context = "..." + context
                    if end < len(line):
                        context = context + "..."

                    hits.append({
                        "file": str(file_path),
                        "filename": file_path.name,
                        "line": line_num,
                        "word": word,
                        "context": context,
                        "position": match.start(),
                    })

        return hits

    def scan_directory(
        self,
        directory: str,
        extensions: Optional[List[str]] = None,
        recursive: bool = True,
    ) -> Dict:
        """批量扫描目录下所有文件"""
        directory = Path(directory)
        if not directory.exists():
            return {"error": f"路径不存在: {directory}"}

        if extensions is None:
            extensions = [".docx", ".md", ".txt", ".text", ".pptx", ".pdf", ".xlsx"]

        ext_set = set(e.lower() if e.startswith(".") else f".{e.lower()}" for e in extensions)
        pattern = "**/*" if recursive else "*"

        all_hits = []
        scanned_files = 0
        hit_files = 0

        for file_path in directory.glob(pattern):
            if not file_path.is_file():
                continue
            if file_path.suffix.lower() not in ext_set:
                continue

            scanned_files += 1
            hits = self.scan_file(str(file_path))

            if hits:
                hit_files += 1
                all_hits.extend(hits)

        # 按文件分组统计
        file_summary = {}
        for hit in all_hits:
            fname = hit["filename"]
            if fname not in file_summary:
                file_summary[fname] = {"file": hit["file"], "hits": [], "word_set": set()}
            file_summary[fname]["hits"].append(hit)
            file_summary[fname]["word_set"].add(hit["word"])

        # 转换为可序列化格式
        summary_list = []
        for fname, info in sorted(file_summary.items(), key=lambda x: -len(x[1]["hits"])):
            summary_list.append({
                "filename": fname,
                "file": info["file"],
                "hit_count": len(info["hits"]),
                "unique_words": sorted(info["word_set"]),
            })

        # 词频统计
        word_freq = {}
        for hit in all_hits:
            w = hit["word"]
            word_freq[w] = word_freq.get(w, 0) + 1

        report = {
            "scan_time": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            "directory": str(directory),
            "total_files_scanned": scanned_files,
            "files_with_hits": hit_files,
            "total_hits": len(all_hits),
            "unique_words_found": list(word_freq.keys()),
            "word_frequency": dict(sorted(word_freq.items(), key=lambda x: -x[1])),
            "file_summary": summary_list,
            "details": all_hits,
        }

        return report

    def save_report(self, report: Dict, output_path: Optional[str] = None) -> str:
        """保存扫描报告为 JSON"""
        if output_path is None:
            self.report_dir.mkdir(parents=True, exist_ok=True)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = str(self.report_dir / f"scan_report_{timestamp}.json")

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        output_path.write_text(
            json.dumps(report, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
        return str(output_path)

    def save_report_csv(self, report: Dict, output_path: Optional[str] = None) -> str:
        """保存扫描报告为 CSV（方便在 Excel 中查看）"""
        import csv

        if output_path is None:
            self.report_dir.mkdir(parents=True, exist_ok=True)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = str(self.report_dir / f"scan_report_{timestamp}.csv")

        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        with open(output_path, "w", encoding="utf-8-sig", newline="") as f:
            writer = csv.writer(f)
            writer.writerow(["文件名", "行号", "敏感词", "上下文", "文件路径"])
            for hit in report.get("details", []):
                writer.writerow([
                    hit["filename"],
                    hit["line"],
                    hit["word"],
                    hit["context"],
                    hit["file"],
                ])

        return str(output_path)

    # ================================================================
    # 文本提取
    # ================================================================

    def _extract_text(self, file_path: Path, ext: str) -> str:
        """从文件中提取纯文本"""
        if ext in (".md", ".txt", ".text"):
            for encoding in ["utf-8", "gbk", "gb2312", "latin-1"]:
                try:
                    return file_path.read_text(encoding=encoding)
                except (UnicodeDecodeError, UnicodeError):
                    continue
            return ""

        if ext == ".docx":
            try:
                from docx import Document
                doc = Document(str(file_path))
                parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
                # 也检查表格
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                parts.append(cell.text)
                return "\n".join(parts)
            except Exception:
                return ""

        if ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(str(file_path))
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            parts.append(shape.text_frame.text)
                return "\n".join(parts)
            except Exception:
                return ""

        if ext == ".pdf":
            try:
                import fitz
                doc = fitz.open(str(file_path))
                parts = [page.get_text("text") for page in doc]
                doc.close()
                return "\n".join(parts)
            except Exception:
                return ""

        if ext in (".xlsx", ".xls"):
            try:
                from openpyxl import load_workbook
                wb = load_workbook(str(file_path), data_only=True)
                parts = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    for row in ws.iter_rows(values_only=True):
                        for cell in row:
                            if cell is not None:
                                parts.append(str(cell))
                wb.close()
                return "\n".join(parts)
            except Exception:
                return ""

        return ""

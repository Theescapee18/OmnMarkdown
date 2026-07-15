"""
OmnMarkdown 批量转换核心模块
支持: .docx, .pdf, .pptx, .xlsx, .txt
支持: 标题、列表、表格、图片提取、加粗/斜体、超链接
"""

import os
import re
import io
import hashlib
from pathlib import Path
from typing import List, Optional, Tuple

try:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph
    from PIL import Image
except ImportError:
    pass

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt as PptxPt
except ImportError:
    Presentation = None

try:
    from openpyxl import load_workbook
except ImportError:
    load_workbook = None


class Word2MarkdownConverter:
    """多格式文档转 Markdown 转换器（.docx, .pdf, .pptx, .xlsx, .txt）"""

    # 支持的格式
    SUPPORTED_EXTS = {".docx", ".pdf", ".pptx", ".xlsx", ".xls", ".txt", ".text", ".md"}

    def __init__(
        self,
        extract_images: bool = True,
        image_format: str = "png",
        image_dir: str = "images",
        heading_offset: int = 0,
        max_image_width: int = 0,
    ):
        self.extract_images = extract_images
        self.image_format = image_format
        self.image_dir = image_dir
        self.heading_offset = heading_offset
        self.max_image_width = max_image_width
        self._image_counter = 0

    def convert_file(self, input_path: str, output_path: Optional[str] = None) -> str:
        """转换单个文件，自动识别格式，返回输出路径"""
        input_path = Path(input_path)
        ext = input_path.suffix.lower()

        if ext not in self.SUPPORTED_EXTS:
            raise ValueError(f"不支持的格式: {ext}，支持: {', '.join(sorted(self.SUPPORTED_EXTS))}")

        if output_path is None:
            output_path = input_path.with_suffix(".md")
        else:
            output_path = Path(output_path)

        output_path.parent.mkdir(parents=True, exist_ok=True)

        # 图片保存到以文档名命名的子目录下，如: output/报告/images/
        self._doc_name = output_path.stem
        self._doc_asset_dir = output_path.parent / self._doc_name

        # 勾选提取图片时，预创建 images 文件夹
        if self.extract_images:
            img_dir = self._doc_asset_dir / self.image_dir
            img_dir.mkdir(parents=True, exist_ok=True)

        # 根据格式路由到对应的转换器
        if ext == ".docx":
            md_content = self._convert_docx(input_path)
        elif ext == ".pdf":
            md_content = self._convert_pdf(input_path)
        elif ext == ".pptx":
            md_content = self._convert_pptx(input_path)
        elif ext in (".xlsx", ".xls"):
            md_content = self._convert_excel(input_path)
        elif ext in (".txt", ".text", ".md"):
            md_content = self._convert_text(input_path)
        else:
            raise ValueError(f"未实现的格式: {ext}")

        output_path.write_text(md_content, encoding="utf-8")
        return str(output_path)

    def _convert_docx(self, input_path: Path) -> str:
        """转换 .docx 文件"""
        doc = Document(str(input_path))
        return self._convert_document(doc, input_path)

    def _convert_document(self, doc: Document, source_path: Path) -> str:
        """转换整个 docx 文档"""
        self._image_counter = 0
        lines = []

        for element in doc.element.body:
            tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

            if tag == "p":
                para = Paragraph(element, doc)
                lines.append(self._convert_paragraph(para, source_path, self._doc_asset_dir))
            elif tag == "tbl":
                table = Table(element, doc)
                lines.append(self._convert_table(table, self._doc_asset_dir))
            elif tag == "sdt":
                # 结构化文档标签（目录等）
                for child_para in element.iter(qn("w:p")):
                    para = Paragraph(child_para, doc)
                    lines.append(self._convert_paragraph(para, source_path, self._doc_asset_dir))

        content = "\n\n".join(line for line in lines if line is not None)
        # 清理多余空行
        content = re.sub(r"\n{3,}", "\n\n", content)
        return content.strip() + "\n"

    def _convert_paragraph(self, para: Paragraph, source_path: Path, output_dir: Path = None) -> str:
        """转换段落为 Markdown"""
        style_name = para.style.name.lower() if para.style else ""

        # 空段落
        if not para.text.strip() and not self._has_images(para):
            return ""

        # 标题
        if style_name.startswith("heading"):
            level = self._extract_heading_level(style_name)
            if level:
                level = min(level + self.heading_offset, 6)
                text = self._convert_runs(para, source_path)
                return f"{'#' * level} {text}"

        # 列表
        if para._element.xpath('.//w:numPr'):
            return self._convert_list_item(para, source_path)

        # 引用块
        if "quote" in style_name or style_name == "intense quote":
            text = self._convert_runs(para, source_path)
            return f"> {text}"

        # 普通段落（含图片）
        text = self._convert_runs(para, source_path)
        images = self._extract_images(para, output_dir or source_path) if self.extract_images else []

        if images:
            img_lines = [f"![{img['alt']}]({img['path']})" for img in images]
            if text.strip():
                return text + "\n\n" + "\n\n".join(img_lines)
            return "\n\n".join(img_lines)

        return text

    def _convert_runs(self, para: Paragraph, source_path: Path) -> str:
        """转换文本 runs，处理加粗/斜体/代码/超链接"""
        parts = []

        for run in para.runs:
            text = run.text
            if not text:
                continue

            # 代码样式
            if run.font.name and "mono" in run.font.name.lower():
                text = f"`{text}`"
            else:
                # 加粗 + 斜体
                if run.bold and run.italic:
                    text = f"***{text}***"
                elif run.bold:
                    text = f"**{text}**"
                elif run.italic:
                    text = f"*{text}*"

                # 下划线
                if run.underline:
                    text = f"<u>{text}</u>"

                # 删除线
                if run.font.strike:
                    text = f"~~{text}~~"

            parts.append(text)

        # 超链接
        for hyperlink in para._element.findall(qn("w:hyperlink")):
            r_id = hyperlink.get(qn("r:id"))
            if r_id and r_id in para.part.rels:
                url = para.part.rels[r_id].target_ref
                link_text = "".join(
                    r.text for r in hyperlink.findall(f".//{qn('w:t')}") if r.text
                )
                if link_text and url:
                    parts.append(f"[{link_text}]({url})")

        return "".join(parts)

    def _extract_heading_level(self, style_name: str) -> Optional[int]:
        """从样式名提取标题级别"""
        match = re.search(r"(\d+)", style_name)
        if match:
            return int(match.group(1))
        if "title" in style_name:
            return 1
        if "subtitle" in style_name:
            return 2
        return None

    def _convert_list_item(self, para: Paragraph, source_path: Path) -> str:
        """转换列表项"""
        text = self._convert_runs(para, source_path)

        # 获取缩进级别
        indent_level = 0
        num_pr = para._element.find(qn("w:pPr"))
        if num_pr is not None:
            ilvl = num_pr.find(qn("w:ilvl"))
            if ilvl is not None:
                indent_level = int(ilvl.get(qn("w:val"), "0"))

        # 判断有序/无序
        is_ordered = False
        num_id_elem = para._element.find(f".//{qn('w:numId')}")
        if num_id_elem is not None:
            num_id = num_id_elem.get(qn("w:val"))
            if num_id and para.part.document.part.numbering_part:
                try:
                    numbering = para.part.document.part.numbering_part.numbering_definitions
                    for abstract_num in numbering.abstract_nums:
                        if str(abstract_num.abstractNumId) == num_id:
                            for lvl in abstract_num.levels:
                                if lvl.numFmt in ("decimal", "upperLetter", "lowerLetter"):
                                    is_ordered = True
                                    break
                except Exception:
                    pass

        prefix = "  " * indent_level
        if is_ordered:
            return f"{prefix}1. {text}"
        return f"{prefix}- {text}"

    def _convert_table(self, table: Table, output_dir: Path = None) -> str:
        """转换表格为 Markdown 表格（支持单元格内图片提取）"""
        rows = table.rows
        if not rows:
            return ""

        extra_images = []
        md_rows = []
        for i, row in enumerate(rows):
            cells = []
            for cell in row.cells:
                cell_parts = []
                for para in cell.paragraphs:
                    # 提取文本
                    text = para.text
                    # 提取单元格内图片
                    if self.extract_images and output_dir:
                        imgs = self._extract_images(para, output_dir)
                        for img in imgs:
                            text += f" ![{img['alt']}]({img['path']})"
                    if text.strip():
                        cell_parts.append(text.strip())
                cell_text = " ".join(cell_parts).replace("\n", " ").replace("|", "\\|")
                cells.append(cell_text)
            md_rows.append("| " + " | ".join(cells) + " |")

            # 在表头后添加分隔行
            if i == 0:
                separator = "| " + " | ".join(["---"] * len(cells)) + " |"
                md_rows.append(separator)

        result = "\n".join(md_rows)
        if extra_images:
            result += "\n\n" + "\n\n".join(
                f"![{img['alt']}]({img['path']})" for img in extra_images
            )
        return result

    def _has_images(self, para: Paragraph) -> bool:
        """检查段落是否包含图片"""
        drawings = para._element.findall(f".//{qn('wp:inline')}")
        drawings += para._element.findall(f".//{qn('wp:anchor')}")
        return len(drawings) > 0

    def _extract_images(self, para: Paragraph, output_dir: Path) -> List[dict]:
        """从段落提取图片，保存到 output_dir 下的 image 子目录"""
        images = []
        drawings = para._element.findall(f".//{qn('wp:inline')}")
        drawings += para._element.findall(f".//{qn('wp:anchor')}")

        for drawing in drawings:
            blip = drawing.find(f".//{qn('a:blip')}")
            if blip is None:
                continue

            embed = blip.get(qn("r:embed"))
            if not embed:
                continue

            try:
                rel = para.part.rels[embed]
                image_data = rel.target_part.blob
                content_type = rel.target_part.content_type

                # 确定扩展名
                ext = self._get_image_extension(content_type)

                # 生成文件名
                self._image_counter += 1
                img_hash = hashlib.md5(image_data).hexdigest()[:8]
                filename = f"img_{self._image_counter:04d}_{img_hash}.{ext}"

                # 保存图片到 output_dir/image/ 下
                img_dir = output_dir / self.image_dir
                img_dir.mkdir(parents=True, exist_ok=True)
                img_path = img_dir / filename
                img_path.write_bytes(image_data)

                # Markdown 中的相对路径（包含文档名子目录）
                rel_path = f"{self._doc_name}/{self.image_dir}/{filename}"
                images.append({"path": rel_path, "alt": f"图片 {self._image_counter}"})

            except Exception as e:
                print(f"  [警告] 提取图片失败: {e}")

        return images

    def _get_image_extension(self, content_type: str) -> str:
        """根据 MIME 类型获取扩展名"""
        mime_map = {
            "image/png": "png",
            "image/jpeg": "jpg",
            "image/gif": "gif",
            "image/bmp": "bmp",
            "image/tiff": "tiff",
            "image/svg+xml": "svg",
            "image/webp": "webp",
        }
        return mime_map.get(content_type, self.image_format)

    def _resize_image(self, image_data: bytes, max_width: int) -> bytes:
        """缩放图片"""
        try:
            img = Image.open(io.BytesIO(image_data))
            if img.width > max_width:
                ratio = max_width / img.width
                new_size = (max_width, int(img.height * ratio))
                img = img.resize(new_size, Image.LANCZOS)
                buf = io.BytesIO()
                fmt = "PNG" if img.mode == "RGBA" else "JPEG"
                img.save(buf, format=fmt, quality=85)
                return buf.getvalue()
        except Exception:
            pass
        return image_data

    # ================================================================
    # PDF 转换
    # ================================================================

    def _convert_pdf(self, input_path: Path) -> str:
        """转换 PDF 文件（使用 PyMuPDF）"""
        if fitz is None:
            raise ImportError("需要安装 PyMuPDF: pip install PyMuPDF")

        doc = fitz.open(str(input_path))
        lines = []
        self._image_counter = 0

        for page_num, page in enumerate(doc, 1):
            # 提取文本
            text = page.get_text("text")
            if text.strip():
                lines.append(text.strip())

            # 提取图片
            if self.extract_images:
                image_list = page.get_images(full=True)
                for img_info in image_list:
                    xref = img_info[0]
                    try:
                        pix = fitz.Pixmap(doc, xref)
                        if pix.n >= 5:  # CMYK 转 RGB
                            pix = fitz.Pixmap(fitz.csRGB, pix)

                        img_data = pix.tobytes("png")
                        self._image_counter += 1
                        img_hash = hashlib.md5(img_data).hexdigest()[:8]
                        filename = f"img_{self._image_counter:04d}_{img_hash}.png"

                        img_dir = self._doc_asset_dir / self.image_dir
                        img_dir.mkdir(parents=True, exist_ok=True)
                        (img_dir / filename).write_bytes(img_data)

                        rel_path = f"{self._doc_name}/{self.image_dir}/{filename}"
                        lines.append(f"![图片 {self._image_counter}]({rel_path})")
                    except Exception as e:
                        print(f"  [警告] PDF 图片提取失败: {e}")

        doc.close()
        return "\n\n".join(lines)

    # ================================================================
    # PPT 转换
    # ================================================================

    def _convert_pptx(self, input_path: Path) -> str:
        """转换 PowerPoint 文件"""
        if Presentation is None:
            raise ImportError("需要安装 python-pptx: pip install python-pptx")

        prs = Presentation(str(input_path))
        lines = []
        self._image_counter = 0

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_lines = []

            for shape in slide.shapes:
                # 文本框
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            # 简单判断层级
                            if para.level == 0:
                                slide_lines.append(text)
                            else:
                                slide_lines.append(f"- {text}")

                # 表格
                if shape.has_table:
                    table = shape.table
                    md_rows = []
                    for i, row in enumerate(table.rows):
                        cells = [cell.text.strip().replace("\n", " ").replace("|", "\\|") for cell in row.cells]
                        md_rows.append("| " + " | ".join(cells) + " |")
                        if i == 0:
                            md_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                    slide_lines.append("\n".join(md_rows))

                # 图片
                if self.extract_images and shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    try:
                        image = shape.image
                        img_data = image.blob
                        ext = image.content_type.split("/")[-1]
                        if ext == "jpeg":
                            ext = "jpg"

                        self._image_counter += 1
                        img_hash = hashlib.md5(img_data).hexdigest()[:8]
                        filename = f"img_{self._image_counter:04d}_{img_hash}.{ext}"

                        img_dir = self._doc_asset_dir / self.image_dir
                        img_dir.mkdir(parents=True, exist_ok=True)
                        (img_dir / filename).write_bytes(img_data)

                        rel_path = f"{self._doc_name}/{self.image_dir}/{filename}"
                        slide_lines.append(f"![图片 {self._image_counter}]({rel_path})")
                    except Exception as e:
                        print(f"  [警告] PPT 图片提取失败: {e}")

            if slide_lines:
                lines.append(f"## 幻灯片 {slide_num}\n")
                lines.extend(slide_lines)

        return "\n\n".join(lines)

    # ================================================================
    # Excel 转换
    # ================================================================

    def _convert_excel(self, input_path: Path) -> str:
        """转换 Excel 文件"""
        if load_workbook is None:
            raise ImportError("需要安装 openpyxl: pip install openpyxl")

        wb = load_workbook(str(input_path), data_only=True)
        lines = []

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"## {sheet_name}\n")

            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                lines.append("（空表）")
                continue

            # 找到最大列数
            max_cols = max(len(row) for row in rows)

            md_rows = []
            for i, row in enumerate(rows):
                cells = []
                for j in range(max_cols):
                    val = row[j] if j < len(row) else ""
                    cell_text = str(val).strip() if val is not None else ""
                    cell_text = cell_text.replace("\n", " ").replace("|", "\\|")
                    cells.append(cell_text)

                md_rows.append("| " + " | ".join(cells) + " |")

                # 第一行作为表头
                if i == 0:
                    md_rows.append("| " + " | ".join(["---"] * max_cols) + " |")

            lines.append("\n".join(md_rows))

        wb.close()
        return "\n\n".join(lines)

    # ================================================================
    # 纯文本转换
    # ================================================================

    def _convert_text(self, input_path: Path) -> str:
        """转换纯文本文件"""
        # 尝试多种编码
        for encoding in ["utf-8", "gbk", "gb2312", "latin-1"]:
            try:
                content = input_path.read_text(encoding=encoding)
                return content.strip() + "\n"
            except (UnicodeDecodeError, UnicodeError):
                continue
        raise ValueError(f"无法解码文件: {input_path.name}")
